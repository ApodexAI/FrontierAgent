"""Best-effort text rendering for benchmark deliverable files."""

from __future__ import annotations

from pathlib import Path

_TEXT_EXTENSIONS = {".csv", ".html", ".json", ".md", ".tsv", ".txt", ".yaml"}


def extract_file_text(path: Path, *, per_file_max: int = 20_000) -> str:
    ext = path.suffix.lower()
    try:
        if ext in _TEXT_EXTENSIONS:
            return path.read_text(encoding="utf-8", errors="replace")[:per_file_max]
        if ext in {".xlsx", ".xlsm"}:
            import openpyxl
            workbook = openpyxl.load_workbook(path, data_only=True, read_only=True)
            rows: list[str] = []
            for sheet in workbook.worksheets:
                rows.append(f"-- sheet: {sheet.title} --")
                rows.extend("\t".join("" if value is None else str(value) for value in row) for row in sheet.iter_rows(values_only=True))
            return "\n".join(rows)[:per_file_max]
        if ext == ".docx":
            import docx
            return "\n".join(p.text for p in docx.Document(str(path)).paragraphs)[:per_file_max]
        if ext == ".pptx":
            from pptx import Presentation
            # has_text_frame is the library's own guard; BaseShape does not
            # declare text_frame.
            parts = [
                shape.text_frame.text  # pyright: ignore[reportAttributeAccessIssue]
                for slide in Presentation(str(path)).slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            ]
            return "\n".join(parts)[:per_file_max]
        if ext == ".pdf":
            from pypdf import PdfReader
            return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)[:per_file_max]
    except Exception as exc:
        return f"[could not extract {path.name}: {exc}]"
    return f"[binary file: {path.name} ({path.stat().st_size} bytes)]"


def render_dir(outputs_dir: Path, *, max_chars: int = 40_000) -> str:
    if not outputs_dir.is_dir():
        return ""
    parts: list[str] = []
    for path in sorted(item for item in outputs_dir.rglob("*") if item.is_file()):
        parts.append(f"\n===== /outputs/{path.relative_to(outputs_dir)} =====")
        parts.append(extract_file_text(path))
    return "\n".join(parts).strip()[:max_chars]

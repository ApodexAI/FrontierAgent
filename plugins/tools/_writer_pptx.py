
# pptx writing: incremental load-modify-save; the anchor is a slide number (1-based) + a placeholder role, or anchor text (cross-deck find/replace).
# Implemented with python-pptx; feature set informed by archipelago's (Apache-2.0) slides_server, addressing and parameters are this project's own design.
#
# Operations (op):
#   create(slides=[{layout?, title?, subtitle?, bullets?[], table?[[]], notes?}])
#   add_slide(layout="title_and_content", title?, bullets?[], index?)   index omitted = at the end
#   set_text(slide, placeholder, text)        slide = 1-based page number; placeholder = title|body|subtitle
#   add_textbox(slide, text, x=1, y=1, w=8, h=1)    inches
#   add_table(slide, rows=[[...]], x=0.5, y=1.5, w=9, h=3, header=True)
#   add_image(slide, image_path, x=1, y=1, w?, h?)
#   replace_text(find, replace, slide?)       slide omitted = the whole deck
#   set_notes(slide, text)
#   add_shape(slide, shape=rectangle|oval|..., x,y,w,h, text?, fill_color?, line_color?)
#   add_chart(slide, chart_type=bar|column|line|pie|..., categories[], series={name:[vals]},
#             x?,y?,w?,h?, title?)
#   format_text(slide, find, bold?, italic?, underline?, font_size?, font_color?)
#   duplicate_slide(slide)                     copy a slide and append it at the end
#   delete_slide(slide)
import os as _os

_LAYOUTS = {"title": 0, "title_and_content": 1, "section_header": 2,
            "two_content": 3, "title_only": 5, "blank": 6}
_PH = {"title": 0, "body": 1, "subtitle": 1, "content": 1}

# Shape name → MSO_SHAPE (the enum is fetched lazily, see _shape_enum)
_SHAPES = {"rectangle": "RECTANGLE", "rounded_rectangle": "ROUNDED_RECTANGLE",
           "oval": "OVAL", "ellipse": "OVAL", "diamond": "DIAMOND",
           "triangle": "ISOCELES_TRIANGLE", "right_arrow": "RIGHT_ARROW",
           "left_arrow": "LEFT_ARROW", "up_arrow": "UP_ARROW", "down_arrow": "DOWN_ARROW",
           "pentagon": "PENTAGON", "chevron": "CHEVRON", "star": "STAR_5_POINT",
           "cloud": "CLOUD", "heart": "HEART"}
# Chart name → XL_CHART_TYPE
_CHARTS = {"bar": "BAR_CLUSTERED", "column": "COLUMN_CLUSTERED", "line": "LINE",
           "line_markers": "LINE_MARKERS", "pie": "PIE", "doughnut": "DOUGHNUT",
           "area": "AREA", "radar": "RADAR"}


def _rgb(c):
    """6-digit hex → RGBColor."""
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(str(c).lstrip("#").upper())


def _pp():
    _ensure("pptx", "python-pptx")
    import pptx
    return pptx


def _slide(prs, n):
    """1-based page number → slide object."""
    idx = int(n) - 1
    if idx < 0 or idx >= len(prs.slides.__iter__.__self__._sldIdLst):
        raise ValueError(f"slide {n} out of range (deck has {len(prs.slides._sldIdLst)} slides)")
    return list(prs.slides)[idx]


def _align_pptx(p, align):
    if not align:
        return
    from pptx.enum.text import PP_ALIGN
    m = {"left": "LEFT", "center": "CENTER", "right": "RIGHT",
         "justify": "JUSTIFY"}.get(str(align).lower())
    if m:
        p.alignment = getattr(PP_ALIGN, m)


def _set_pptx_ea(run, font):
    """Add a:ea (the East Asian font) to a pptx run — font.name only writes a:latin, and CJK characters read a:ea."""
    try:
        from pptx.oxml.ns import qn as _q
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:ea", "a:cs"):
            el = rPr.find(_q(tag))
            if el is None:
                el = rPr.makeelement(_q(tag), {})
                rPr.append(el)
            el.set("typeface", font)
    except Exception:
        pass


def _rt_para(p, text):
    """Write RichText (normalised by _norm_runs) run by run into one pptx paragraph p. A link becomes a real hyperlink."""
    from pptx.util import Pt
    for rd in _norm_runs(text):
        r = p.add_run()
        r.text = rd.get("text", "")
        f = r.font
        if rd.get("bold") is not None:
            f.bold = bool(rd["bold"])
        if rd.get("italic") is not None:
            f.italic = bool(rd["italic"])
        if rd.get("underline") is not None:
            f.underline = bool(rd["underline"])
        if rd.get("size"):
            f.size = Pt(float(rd["size"]))
        if rd.get("color"):
            f.color.rgb = _rgb(rd["color"])
        if rd.get("font"):
            f.name = rd["font"]
            # python-pptx's font.name only writes a:latin — CJK characters use a:ea, so
            # without setting it the requested
            # font has no effect on Chinese/Japanese/Korean text (the same problem as docx's w:eastAsia).
            _set_pptx_ea(r, rd["font"])
        if rd.get("link"):
            r.hyperlink.address = rd["link"]


def _no_bullet(p):
    """Turn off a paragraph's bullet (already-numbered content, or content that needs no dot): insert a:buNone into pPr."""
    from pptx.oxml.ns import qn
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def _set_autofit(tf, mode):
    """Text-box autofit: none / shrink_text (shrink the font) / resize_shape (grow the shape)."""
    if not mode:
        return
    from pptx.enum.text import MSO_AUTO_SIZE
    tf.word_wrap = True
    m = str(mode).lower()
    if m in ("shrink", "shrink_text"):
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    elif m in ("resize", "resize_shape"):
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    elif m == "none":
        tf.auto_size = MSO_AUTO_SIZE.NONE


def _find_ph(slide, role):
    """Find a placeholder by its semantic type (fixes: body and subtitle both map to idx 1 in _PH,
    so writing both on one slide clobbered each other and silently lost content). When the exact
    type is missing it falls back to idx, but skips placeholders belonging to another
    role's dedicated type (no stealing); still nothing → None (the caller falls back to a text box)."""
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER as _P
        WANT = {"title": {_P.TITLE, _P.CENTER_TITLE},
                "subtitle": {_P.SUBTITLE},
                "body": {_P.BODY, _P.OBJECT},
                "content": {_P.BODY, _P.OBJECT}}
        AVOID = {"body": {_P.SUBTITLE, _P.TITLE, _P.CENTER_TITLE},
                 "content": {_P.SUBTITLE, _P.TITLE, _P.CENTER_TITLE},
                 "subtitle": {_P.BODY, _P.OBJECT, _P.TITLE, _P.CENTER_TITLE}}
        want, avoid = WANT.get(role), AVOID.get(role, set())
    except Exception:  # enum unavailable → fall back to matching on idx alone
        want, avoid = None, set()
    phs = list(slide.placeholders)
    if want:
        for ph in phs:
            try:
                if ph.placeholder_format.type in want:
                    return ph
            except Exception:
                continue
    idx = _PH.get(role, 1)
    for ph in phs:
        try:
            if ph.placeholder_format.idx == idx and ph.placeholder_format.type not in avoid:
                return ph
        except Exception:
            continue
    return None


def _set_ph(slide, role, text):
    """Write RichText into a placeholder. role = title/body/subtitle. A missing placeholder falls back to a text box."""
    from pptx.util import Inches
    tf = None
    if role == "title" and slide.shapes.title is not None:
        tf = slide.shapes.title.text_frame
    else:
        ph = _find_ph(slide, role)
        if ph is not None:
            tf = ph.text_frame
    if tf is None:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8), Inches(1))
        tf = tb.text_frame
    tf.clear()
    _rt_para(tf.paragraphs[0], text)
    return True


def _add_body(slide, body):
    """Structured list for the body placeholder (cures flattening and forced bullets). body may be:
    - a list ['a','b'] (the old bullets form, still supported, defaults to level 0 with a dot);
    - dict {items:[{text:RichText, level:0-4, bullet:true|false}], autofit?}。"""
    from pptx.util import Inches
    if isinstance(body, dict):
        items = body.get("items", [])
        autofit = body.get("autofit")
    else:
        items = list(body or [])
        autofit = None
    ph = _find_ph(slide, "body")
    if ph is None:
        ph = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(4))
    tf = ph.text_frame
    tf.clear()
    for i, it in enumerate(items):
        if isinstance(it, str):
            it = {"text": it}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = int(it.get("level", 0))
        _rt_para(p, it.get("text", ""))
        if it.get("bullet") is False:
            _no_bullet(p)
    _set_autofit(tf, autofit)


def _add_table(slide, rows, x, y, w, h, header):
    from pptx.util import Inches
    nr, nc = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            tbl.cell(i, j).text = str(val)
    return gf


def _build_slide(prs, sd):
    layout = prs.slide_layouts[_LAYOUTS.get(sd.get("layout", "title_and_content"), 1)]
    slide = prs.slides.add_slide(layout)
    if sd.get("title"):
        _set_ph(slide, "title", sd["title"])
    if sd.get("subtitle"):
        _set_ph(slide, "subtitle", sd["subtitle"])
    if sd.get("body") is not None:
        _add_body(slide, sd["body"])
    elif sd.get("bullets"):
        _add_body(slide, sd["bullets"])  # the old bullets form
    if sd.get("table"):
        _add_table(slide, sd["table"], 0.5, 2.0, 9, 3, True)
    if sd.get("notes"):
        ntf = slide.notes_slide.notes_text_frame
        ntf.clear()
        _rt_para(ntf.paragraphs[0], sd["notes"])
    return slide


def _pptx_write(path, op, args):
    pptx = _pp()

    if op == "create":
        if _os.path.exists(path) and not args.get("overwrite"):
            return _res(f"create refused: {path} already exists — use add_slide/set_text to "
                        "edit, or pass overwrite:true to rebuild", ok=False)
        slides = args.get("slides", [])
        if not slides:
            return _res(f"create wrote nothing to {path}: no slides given", ok=False)
        prs = pptx.Presentation()
        for sd in slides:
            _build_slide(prs, sd)
        _os.makedirs(_os.path.dirname(path) or ".", exist_ok=True)
        prs.save(path)
        return _res(f"created pptx: {path}", counts={"slide": len(slides)})

    if not _os.path.exists(path):
        return f"[error] file not found (edit needs existing file): {path}"
    prs = pptx.Presentation(path)

    if op == "add_slide":
        sd = {k: args[k] for k in ("layout", "title", "subtitle", "body", "bullets", "notes", "table")
              if k in args}
        idx = args.get("index")
        slide = _build_slide(prs, sd)
        if idx is not None:  # move to the requested position (1-based)
            lst = prs.slides._sldIdLst
            el = lst[-1]
            lst.remove(el)
            lst.insert(min(int(idx) - 1, len(lst)), el)
        prs.save(path)
        return f"added slide ({sd.get('layout', 'title_and_content')})"

    if op == "set_text":
        slide = _slide(prs, args["slide"])
        _set_ph(slide, args.get("placeholder", "body"), args["text"])
        prs.save(path)
        return f"set {args.get('placeholder', 'body')} on slide {args['slide']}"

    if op == "add_textbox":
        from pptx.util import Inches
        slide = _slide(prs, args["slide"])
        tb = slide.shapes.add_textbox(Inches(args.get("x", 1)), Inches(args.get("y", 1)),
                                      Inches(args.get("w", 8)), Inches(args.get("h", 1)))
        tf = tb.text_frame
        tf.clear()
        _rt_para(tf.paragraphs[0], args.get("text", ""))
        _align_pptx(tf.paragraphs[0], args.get("align_h"))
        _set_autofit(tf, args.get("autofit"))
        prs.save(path)
        return f"added textbox on slide {args['slide']}"

    if op == "add_table":
        slide = _slide(prs, args["slide"])
        _add_table(slide, args["rows"], args.get("x", 0.5), args.get("y", 1.5),
                   args.get("w", 9), args.get("h", 3), args.get("header", True))
        prs.save(path)
        return f"added table on slide {args['slide']}"

    if op == "add_image":
        from pptx.util import Inches
        slide = _slide(prs, args["slide"])
        kw = {}
        if args.get("w"):
            kw["width"] = Inches(args["w"])
        if args.get("h"):
            kw["height"] = Inches(args["h"])
        slide.shapes.add_picture(args["image_path"], Inches(args.get("x", 1)),
                                 Inches(args.get("y", 1)), **kw)
        prs.save(path)
        return f"added image on slide {args['slide']}"

    if op == "set_notes":
        slide = _slide(prs, args["slide"])
        ntf = slide.notes_slide.notes_text_frame
        ntf.clear()
        _rt_para(ntf.paragraphs[0], args["text"])
        prs.save(path)
        return f"set notes on slide {args['slide']}"

    if op == "replace_text":
        find, repl = args["find"], args.get("replace", "")
        target = [_slide(prs, args["slide"])] if args.get("slide") else list(prs.slides)
        cnt = 0
        for slide in target:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text and find in r.text:
                            r.text = r.text.replace(find, repl)
                            cnt += 1
        prs.save(path)
        return _res(f"replaced {cnt} run(s) containing {find!r}",
                    warn=(f"0 matches for {find!r}" if cnt == 0 else None))

    if op == "add_shape":
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
        slide = _slide(prs, args["slide"])
        name = str(args.get("shape", "rectangle")).lower()
        mso = getattr(MSO_SHAPE, _SHAPES.get(name, "RECTANGLE"))
        sp = slide.shapes.add_shape(mso, Inches(args.get("x", 1)), Inches(args.get("y", 1)),
                                    Inches(args.get("w", 2)), Inches(args.get("h", 1)))
        if args.get("fill_color"):
            sp.fill.solid()
            sp.fill.fore_color.rgb = _rgb(args["fill_color"])
        if args.get("line_color"):
            sp.line.color.rgb = _rgb(args["line_color"])
        if args.get("text"):
            sp.text_frame.text = args["text"]
            if args.get("font_color"):
                sp.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(args["font_color"])
            if args.get("font_size"):
                sp.text_frame.paragraphs[0].runs[0].font.size = Pt(float(args["font_size"]))
        prs.save(path)
        return f"added {name} on slide {args['slide']}"

    if op == "add_chart":
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
        slide = _slide(prs, args["slide"])
        ctype = getattr(XL_CHART_TYPE, _CHARTS.get(
            str(args.get("chart_type", "column")).lower(), "COLUMN_CLUSTERED"))
        data = CategoryChartData()
        data.categories = args.get("categories") or []
        series = args.get("series") or {}
        for sname, vals in series.items():
            data.add_series(sname, [float(v) for v in vals])
        gf = slide.shapes.add_chart(
            ctype, Inches(args.get("x", 1)), Inches(args.get("y", 1.5)),
            Inches(args.get("w", 8)), Inches(args.get("h", 4.5)), data)
        if args.get("title"):
            gf.chart.has_title = True
            gf.chart.chart_title.text_frame.text = args["title"]
        prs.save(path)
        return f"added {args.get('chart_type', 'column')} chart on slide {args['slide']}"

    if op == "format_text":
        from pptx.util import Pt
        slide = _slide(prs, args["slide"])
        find = args["find"]
        hits = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if not r.text or find not in r.text:
                        continue
                    f = r.font
                    if args.get("bold") is not None:
                        f.bold = bool(args["bold"])
                    if args.get("italic") is not None:
                        f.italic = bool(args["italic"])
                    if args.get("underline") is not None:
                        f.underline = bool(args["underline"])
                    if args.get("strike") is not None:  # python-pptx has no font.strike, so go through rPr
                        rPr = r._r.get_or_add_rPr()
                        rPr.set("strike", "sngStrike" if args["strike"] else "noStrike")
                    sz = args.get("font_size") or args.get("size")  # size = the parameter name the documented contract uses
                    if sz:
                        f.size = Pt(float(sz))
                    col = args.get("font_color") or args.get("color")
                    if col:
                        f.color.rgb = _rgb(col)
                    hits += 1
        prs.save(path)
        return _res(f"formatted {hits} run(s) matching {find!r} on slide {args['slide']}",
                    warn=(f"0 runs matched {find!r}" if hits == 0 else None))

    if op == "duplicate_slide":
        import copy as _copy
        src = _slide(prs, args["slide"])
        layout = src.slide_layout
        new = prs.slides.add_slide(layout)
        # Drop the placeholders the layout brings and deep-copy the source slide's shapes instead
        for sh in list(new.shapes):
            sh._element.getparent().remove(sh._element)
        for sh in src.shapes:
            new.shapes._spTree.append(_copy.deepcopy(sh._element))
        # Also copy the source slide's relationships (rels for images / charts / hyperlinks)
        # and rewrite the rIds inside the copied XML —
        # copying only the shape XML leaves r:embed pointing at a relationship the new slide does not have, losing assets and making PowerPoint report corruption.
        _RNS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        id_map = {}
        for rid, rel in list(src.part.rels.items()):
            if rel.reltype.endswith("/slideLayout") or rel.reltype.endswith("/notesSlide"):
                continue  # the new slide already has its own layout; notes do not travel with shapes
            if rel.is_external:
                new_rid = new.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                new_rid = new.part.relate_to(rel.target_part, rel.reltype)
            if new_rid != rid:
                id_map[rid] = new_rid
        if id_map:
            for el in new.shapes._spTree.iter():
                for attr, val in list(el.attrib.items()):
                    if attr.startswith("{" + _RNS + "}") and val in id_map:
                        el.set(attr, id_map[val])
        prs.save(path)
        return f"duplicated slide {args['slide']} -> slide {len(prs.slides._sldIdLst)}"

    if op == "delete_slide":
        idx = int(args["slide"]) - 1
        lst = prs.slides._sldIdLst
        ids = list(lst)
        if idx < 0 or idx >= len(ids):
            return f"[error] slide {args['slide']} out of range"
        lst.remove(ids[idx])
        prs.save(path)
        return f"deleted slide {args['slide']}"

    if op == "set_slide_size":
        # Slide canvas size / aspect ratio (the most common pptx patch in evaluation: switching to widescreen when the template is not 16:9).
        from pptx.util import Inches
        preset = str(args.get("preset") or args.get("aspect") or "").lower().replace("x", ":")
        presets = {"16:9": (13.333, 7.5), "widescreen": (13.333, 7.5),
                   "4:3": (10, 7.5), "standard": (10, 7.5), "16:10": (10, 6.25)}
        if preset in presets:
            w, h = presets[preset]
        else:
            w = float(args.get("width_in", 13.333))
            h = float(args.get("height_in", 7.5))
        prs.slide_width = Inches(w)
        prs.slide_height = Inches(h)
        prs.save(path)
        return f"set slide size {w}in x {h}in ({preset or 'custom'})"

    return f"[error] unknown pptx op: {op}"

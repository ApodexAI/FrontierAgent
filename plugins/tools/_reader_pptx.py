
# pptx parsing: inline code `…` marks parser-added metadata (position / type / format / chart markers), not slide text.
_GROUP_MIN_COMMON = 3  # text shapes sharing this many parameters (backtick bits) collapse into one group, with the shared parameters hoisted to the group header
_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_NS_P_CNVPR = "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
_NS_A_CNVPR = "{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr"


def _emu_in(v):
    """EMU → inches (1 inch = 914400 EMU), 2 decimal places."""
    return round(v / 914400, 2) if v is not None else None


def _run_color(run) -> str:
    try:
        c = run.font.color
        if c is not None and c.type is not None:
            return str(c.rgb)
    except Exception:
        pass
    return ""


def _run_spc(run):
    """Run character spacing a:rPr@spc (units of 1/100 pt)."""
    try:
        rPr = run._r.find(f"{_A}rPr")
        if rPr is not None and rPr.get("spc"):
            return round(int(rPr.get("spc")) / 100, 1)
    except Exception:
        pass
    return None


def _fmt_num(x) -> str:
    """28.0 → '28', 13.5 → '13.5' (drop the pointless .0)."""
    return f"{x:g}"


def _run_md(run, base_size=None, base_color=None) -> str:
    """run → markdown. **bold** / *italic* / <u>underline</u> inline; the trailing
    `size,colour,spacing` note (inline code) records **only deviations**:
    size only when it differs from the shape's main size; colour only when non-black
    (#000000) and different from the shape's main colour; spacing only when set.
    The shape's main size / colour are computed by the caller and noted once at shape
    level, so runs do not repeat default values."""
    t = run.text
    if not t:
        return ""
    lead, trail, core = t[: len(t) - len(t.lstrip())], t[len(t.rstrip()):], t.strip()
    if not core:
        return t
    f = run.font
    if f.bold and f.italic:
        core = f"***{core}***"
    elif f.bold:
        core = f"**{core}**"
    elif f.italic:
        core = f"*{core}*"
    if f.underline:
        core = f"<u>{core}</u>"
    ann = []
    try:
        sz = round(f.size.pt, 1) if f.size is not None else None
    except Exception:
        sz = None
    if sz is not None and sz != base_size:
        ann.append(f"{_fmt_num(sz)}pt")
    col = _run_color(run)
    if col and col != "000000" and col != base_color:
        ann.append(f"#{col}")
    spc = _run_spc(run)
    if spc:
        ann.append(f"spc{_fmt_num(spc)}")
    if ann:
        core = f"{core}`{','.join(ann)}`"
    return lead + core + trail


def _shape_base_fmt(tf):
    """A shape's main size / colour (the most frequent non-empty value; colour ignores black). Used to hoist shared formatting to shape level and note it once."""
    from collections import Counter
    sizes, colors = [], []
    for p in tf.paragraphs:
        for r in p.runs:
            if not r.text.strip():
                continue
            try:
                if r.font.size is not None:
                    sizes.append(round(r.font.size.pt, 1))
            except Exception:
                pass
            colors.append(_run_color(r) or "000000")  # unset colour counts as the default black
    base_size = Counter(sizes).most_common(1)[0][0] if sizes else None
    # The main colour must be a real majority (black included in the count); a black
    # majority is not hoisted (black = default, not annotated).
    # Only hoist when most runs share a non-black colour (e.g. an all-white title), so a
    # single coloured run is not mistaken for a global colour.
    base_color = Counter(colors).most_common(1)[0][0] if colors else None
    if base_color == "000000":
        base_color = None
    return base_size, base_color


def _para_list_kind(p):
    """Paragraph bullet type: ('auto',type,startAt)=numbered / ('char',symbol)=bullet /
    ('none',)=explicitly none / None=inherited (decided by the placeholder default). Reads buAutoNum/buChar/buNone under a:pPr."""
    try:
        pPr = p._p.find(f"{_A}pPr")
    except Exception:
        return None
    if pPr is None:
        return None
    if pPr.find(f"{_A}buNone") is not None:
        return ("none",)
    auto = pPr.find(f"{_A}buAutoNum")
    if auto is not None:
        return ("auto", auto.get("type", "arabicPeriod"), auto.get("startAt"))
    ch = pPr.find(f"{_A}buChar")
    if ch is not None:
        return ("char", ch.get("char", "•"))
    return None


def _ph_type(shape) -> str:
    if not shape.is_placeholder:
        return ""
    try:
        return str(shape.placeholder_format.type).split()[0]
    except Exception:
        return ""


def _shape_fill(shape) -> str:
    """Shape fill colour (solid fill only) → #RRGGBB."""
    try:
        fill = shape.fill
        if fill.type is not None and int(fill.type) == 1:  # MSO_FILL.SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return ""


def _shape_bits(shape, with_size=False) -> list:
    """Shape-level meta bits — **only what carries signal**: position @x,y (to judge
    left/right and relative layout), placeholder type,
    geometry (only when not a rectangle, e.g. arrow / ellipse), fill colour (only when non-black).
    Zero-information items are dropped by default: AUTO_SHAPE type names, RECTANGLE
    geometry, auto-generated names ("Text 3" / "Shape 5"), text-box dimensions.
    with_size=True adds shape dimensions (needed only for text-free visual elements and
    pictures, not for text boxes)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    bits = []
    try:
        if shape.left is not None:
            bits.append(f"@{_emu_in(shape.left)},{_emu_in(shape.top)}")
    except Exception:
        pass
    if with_size:
        try:
            if shape.width is not None:
                bits.append(f"{_emu_in(shape.width)}×{_emu_in(shape.height)}in")
        except Exception:
            pass
    ph = _ph_type(shape)
    if ph:
        bits.append(f"ph={ph}")
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.auto_shape_type is not None:
            g = str(shape.auto_shape_type).split()[0]
            if g != "RECTANGLE":
                bits.append(g)
    except Exception:
        pass
    fill = _shape_fill(shape)
    if fill and fill != "000000":
        bits.append(f"fill:#{fill}")
    return bits


def _text_vert(shape) -> str:
    """Vertical / rotated text direction (a:bodyPr@vert). Horizontal returns "";
    otherwise the direction value, e.g. vert / vert270 / eaVert.
    The content itself is read as usual (direction-independent); this only adds a
    direction annotation."""
    try:
        bodyPr = shape.text_frame._txBody.find(f"{_A}bodyPr")
        if bodyPr is not None:
            v = bodyPr.get("vert")
            if v and v != "horz":
                return v
    except Exception:
        pass
    return ""


def _pptx_text_md(shape):
    """Shape with text → (fmt_bits, lines). fmt_bits = shape-level shared formatting
    (main size / colour, noted once).
    TITLE → ### heading; everything else renders per paragraph, with bullets / numbering
    decided by a:buChar / a:buAutoNum
    (not guessed from the placeholder, which used to lose level-0 bullets and numbering).
    When inherited, a body placeholder defaults to a bullet and
    a free text box at level 0 is a plain line. Run level only records deviations from the
    shape's main size / colour (see _run_md)."""
    tf = shape.text_frame
    base_size, base_color = _shape_base_fmt(tf)
    fmt_bits = []
    if base_size is not None:
        fmt_bits.append(f"{_fmt_num(base_size)}pt")
    if base_color:
        fmt_bits.append(f"#{base_color}")
    vert = _text_vert(shape)  # vertical-direction annotation (nothing for horizontal)
    if vert:
        fmt_bits.append(vert)
    ph = _ph_type(shape)
    is_title = "TITLE" in ph
    is_body = any(k in ph for k in ("BODY", "OBJECT", "SUBTITLE"))
    lines, counters = [], {}
    for p in tf.paragraphs:
        md = "".join(_run_md(r, base_size, base_color) for r in p.runs).strip() or p.text.strip()
        if not md:
            continue
        if is_title:
            lines.append(f"### {md}")
            continue
        lvl = p.level
        indent = "  " * lvl
        kind = _para_list_kind(p)
        if kind and kind[0] == "auto":
            start = int(kind[2]) if len(kind) > 2 and kind[2] else 1
            n = counters.get(lvl, start)
            counters[lvl] = n + 1
            lines.append(f"{indent}{n}. {md}")
        elif kind and kind[0] == "char":
            lines.append(f"{indent}- {md}")
        elif kind and kind[0] == "none":
            lines.append(f"{indent}{md}")
        elif is_body or lvl > 0:  # inherited
            lines.append(f"{indent}- {md}")
        else:
            lines.append(md)
    return fmt_bits, lines


def _pptx_table_md(shape):
    """Table → (rows, cols, markdown rows). Merged cells: a spanned cell (is_spanned)
    repeats the text, so it is left empty and
    only the merge origin (top-left) is kept. The row x col size goes into the caller's `…` meta (marked as added, not original)."""
    tbl = shape.table
    n_rows, n_cols = len(tbl.rows), len(tbl.columns)
    lines = []
    for r in range(n_rows):
        cells = []
        for c in range(n_cols):
            cell = tbl.cell(r, c)
            txt = "" if cell.is_spanned else cell.text.replace("\n", " ").replace("|", "\\|")
            cells.append(txt)
        lines.append("| " + " | ".join(cells) + " |")
        if r == 0:
            lines.append("| " + " | ".join("---" for _ in range(n_cols)) + " |")
    return n_rows, n_cols, lines


def _chart_custom_labels(ch):
    """Custom data-label text (present only when it overrides the number shown on the
    chart; an automatic label = the series value, already in the data table).
    Scans every plot (a combo chart has several)."""
    found = {}
    try:
        for plot in ch.plots:
            for s in plot.series:
                for idx, pt in enumerate(s.points):
                    try:
                        dl = pt.data_label
                        if dl.has_text_frame and dl.text_frame.text.strip():
                            found.setdefault(str(s.name), {})[idx] = dl.text_frame.text.strip()
                    except Exception:
                        pass
    except Exception:
        pass
    return found or None


def _pptx_chart_md(shape):
    """Chart → (bits, lines). bits go into the shape meta (chart type / title / whether
    data labels are shown).
    The real values live in chart1.xml (text extraction cannot reach them), so this goes
    through the chart API for categories x series;
    custom data-label text is read out too (an automatic label equals the series value,
    already in the data table, so it is not repeated)."""
    ch = shape.chart
    bits = ["chart", str(ch.chart_type).split()[0] if ch.chart_type is not None else "?"]
    if ch.has_title:
        bits.append(f'title="{ch.chart_title.text_frame.text}"')
    try:
        if ch.plots[0].has_data_labels:
            bits.append("data-labels-shown")
    except Exception:
        pass
    lines = []
    try:
        # Combo chart = several plots: collect every series of every plot; the category axis comes from the first non-empty plot.
        series, cats = [], []
        for plot in ch.plots:
            try:
                pc = list(plot.categories)
                if pc and not cats:
                    cats = pc
            except Exception:
                pass
            for s in plot.series:
                series.append((str(s.name), list(s.values)))
        if cats:  # has a category axis: category x series table
            header = ["category"] + [s[0] for s in series]
            lines.append("| " + " | ".join(header) + " |")
            lines.append("| " + " | ".join("---" for _ in header) + " |")
            for i, cat in enumerate(cats):
                row = [str(cat)] + [str(s[1][i]) if i < len(s[1]) else "" for s in series]
                lines.append("| " + " | ".join(row) + " |")
        elif series:  # no category axis (scatter / bubble): just list each series' values
            for name, vals in series:
                lines.append(f'- series "{name}": ' + ", ".join(str(v) for v in vals))
        custom = _chart_custom_labels(ch)
        if custom:
            lines.append(f"`custom data-label text: {custom}`")
    except Exception as e:
        lines.append(f"`chart data unavailable: {type(e).__name__}: {e}`")
    return bits, lines


def _pptx_pic_alt(shape) -> str:
    """A picture is opaque to text extraction, so alt text is the only textual description (often falling back to the filename when unset)."""
    try:
        # An lxml element with no children is falsy, so `a or b` is wrong here — the check must be explicit `is None`.
        cNvPr = shape._element.find(f".//{_NS_P_CNVPR}")
        if cNvPr is None:
            cNvPr = shape._element.find(f".//{_NS_A_CNVPR}")
        if cNvPr is not None:
            return cNvPr.get("descr") or ""
    except Exception:
        pass
    return ""


def _tag(bits) -> str:
    """Join meta bits into inline code `…` (parser-added, not original text); emit no tag when empty."""
    return "`" + " ".join(bits) + "`" if bits else ""


# ---- Shape-level geometry: merge "background box + label", detect arrows / connectors (makes flow and relationship diagrams readable) ----

def _bbox(shape):
    """Shape bounding box (left, top, width, height) in inches. None when unavailable."""
    try:
        if shape.left is None:
            return None
        return (_emu_in(shape.left), _emu_in(shape.top),
                _emu_in(shape.width), _emu_in(shape.height))
    except Exception:
        return None


def _shape_text(shape) -> str:
    try:
        return shape.text_frame.text.strip() if getattr(shape, "has_text_frame", False) else ""
    except Exception:
        return ""


def _is_line(shape) -> bool:
    """A connector is treated as "an arrow joining two blocks" (we want the connection,
    not the physical shape):
    connector type / degenerate to zero width-height / auto_shape 'line' preset /
    **block arrows (RIGHT_ARROW etc.) with no text**.
    A block arrow that does have text is treated as ordinary text (not True here)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return True
    except Exception:
        pass
    b = _bbox(shape)
    if b and (b[2] <= 0.05 or b[3] <= 0.05):
        return True
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                ast = str(shape.auto_shape_type).upper()
            except Exception as e:
                return "line" in str(e).lower()  # 'line' preset: accessing auto_shape_type raises
            if "ARROW" in ast and not _shape_text(shape):  # text-free block arrow = connector
                return True
    except Exception:
        pass
    return False


def _is_box(shape) -> bool:
    """An auto shape usable as a "background box" (rectangle / rounded rectangle / ellipse …), not a line."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and not _is_line(shape)
    except Exception:
        return False


def _pair_label_boxes(shapes, bbox):
    """Pair up "text-free background box + the text label inside it" (authors commonly
    draw a colour block and its text as two overlapping shapes).
    Returns {text shape index: background box index} plus the set of indices used as
    backgrounds (skipped on output, already folded into the text line).
    The background box must fully contain the text box (all four sides, small tolerance),
    and the tightest fit (smallest area) wins."""
    TOL = 0.2
    containers = [i for i, s in enumerate(shapes)
                  if _is_box(s) and not _shape_text(s) and bbox[i]]
    text_to_box, used = {}, set()
    for i, s in enumerate(shapes):
        if not _shape_text(s) or not bbox[i]:
            continue
        tl, tt, tw, th = bbox[i]
        best, best_area = None, None
        for ci in containers:
            if ci in used:
                continue
            cl, ct, cw, ch = bbox[ci]
            if (cl - TOL <= tl and ct - TOL <= tt
                    and cl + cw + TOL >= tl + tw and ct + ch + TOL >= tt + th
                    and cw * ch <= 10 * max(tw * th, 0.01)):
                area = cw * ch
                if best is None or area < best_area:
                    best, best_area = ci, area
        if best is not None:
            text_to_box[i] = best
            used.add(best)
    return text_to_box, used


def _label_at_edge(point, edge, shapes, bbox):
    """Find a shape with text whose given edge sits against ``point`` and return its text
    (truncated). Used to work out what an arrow connects.
    edge='right': that shape's right edge ≈ point.x (at the arrow's left end = source);
    'left': its left edge ≈ x (right end = target);
    'bottom' / 'top' likewise for vertical arrows."""
    px, py = point
    for i, s in enumerate(shapes):
        b = bbox[i]
        txt = _shape_text(s)
        if not b or not txt:
            continue
        x0, t, w, h = b
        if edge in ("right", "left") and not (t - 0.2 <= py <= t + h + 0.2):
            continue
        if edge in ("top", "bottom") and not (x0 - 0.2 <= px <= x0 + w + 0.2):
            continue
        near = {"right": abs((x0 + w) - px), "left": abs(x0 - px),
                "bottom": abs((t + h) - py), "top": abs(t - py)}[edge]
        if near <= 0.35:
            return txt[:24]
    return None


def _arrow_info(shape, idx, shapes, bbox) -> dict:
    """Line / arrow → {line, src, dst, arrow, connected}. Resolves direction and tries to identify the two labels it joins."""
    b = bbox[idx]
    if not b:
        return {"line": _tag(["line", "vlm"]), "src": None, "dst": None, "connected": False}
    x0, t, w, h = b
    at = ""
    try:
        at = str(shape.auto_shape_type).upper()  # block arrows (RIGHT_ARROW…) get their direction from this
    except Exception:
        at = ""
    left_end, right_end = (x0, t + h / 2), (x0 + w, t + h / 2)
    top_end, bot_end = (x0 + w / 2, t), (x0 + w / 2, t + h)
    if w >= h and not ("UP_ARROW" in at or "DOWN_ARROW" in at):  # horizontal
        if "LEFT_ARROW" in at and "RIGHT" not in at:  # points left → right block flows to left block
            src = _label_at_edge(right_end, "left", shapes, bbox)
            dst = _label_at_edge(left_end, "right", shapes, bbox)
            arrow = "←"
        else:  # a thin line / connector with no direction defaults to left → right
            src = _label_at_edge(left_end, "right", shapes, bbox)
            dst = _label_at_edge(right_end, "left", shapes, bbox)
            arrow = "→"
    else:  # vertical
        if "UP_ARROW" in at:  # points up → lower block flows to upper block
            src = _label_at_edge(bot_end, "top", shapes, bbox)
            dst = _label_at_edge(top_end, "bottom", shapes, bbox)
            arrow = "↑"
        else:
            src = _label_at_edge(top_end, "bottom", shapes, bbox)
            dst = _label_at_edge(bot_end, "top", shapes, bbox)
            arrow = "↓"
    bits = [f"@{x0},{t}", f"arrow {arrow}"]
    connected = bool(src and dst)
    if src or dst:
        bits.append(f'connects "{src or "?"}"{arrow}"{dst or "?"}"')
    else:
        bits.append("vlm")
    return {"line": _tag(bits), "src": src, "dst": dst, "arrow": arrow, "connected": connected}


def _order_flow(edges):
    """Chain the (src,dst) edge set into 'A → B → C'; non-linear graphs degrade to listing edges as 'A→B, C→D'."""
    if not edges:
        return None
    succ = {s: d for s, d in edges}
    preds = {d for _, d in edges}
    starts = [s for s, _ in edges if s not in preds]
    nodes = {s for s, _ in edges} | {d for _, d in edges}
    if len(starts) == 1:
        chain, seen = [starts[0]], {starts[0]}
        while chain[-1] in succ and succ[chain[-1]] not in seen:
            chain.append(succ[chain[-1]])
            seen.add(chain[-1])
        if len(chain) == len(nodes):
            return " → ".join(chain)
    return ", ".join(f"{s}→{d}" for s, d in edges)


def _pptx_slide_lines(shapes, bbox, text_to_box, used, MSO):
    """Render one slide's content lines (without '<!-- slide N -->' or notes). Returns (lines, needs_vlm).
    Tiered compression: text shapes whose style signature (geometry / size / fill / font
    size and colour, position excluded) is identical and that number >= 2 collapse into a group —
    the shared part becomes the group header and members keep only @position + label;
    an arrow with both ends in the same group folds into that group's 'flow:' line."""
    from collections import defaultdict

    records = []
    for idx, shape in enumerate(shapes):
        if idx in used:
            continue
        rec = {"idx": idx, "kind": "block", "lines": []}
        try:
            if getattr(shape, "has_table", False):
                nr, nc, tbl = _pptx_table_md(shape)
                rec["lines"] = [_tag([*_shape_bits(shape), f"table {nr}×{nc}"]), *tbl]
            elif getattr(shape, "has_chart", False):
                cbits, clines = _pptx_chart_md(shape)
                rec["lines"] = [_tag(_shape_bits(shape) + cbits), *clines]
            elif shape.shape_type == MSO.PICTURE:
                alt = _pptx_pic_alt(shape)
                bits = [*_shape_bits(shape, with_size=True), "image", "vlm"]
                bits.append(f'alt:"{alt}"' if alt else "no-alt")  # alt folded into this line, removing any ambiguity about what it belongs to
                rec["lines"] = [_tag(bits)]
                rec["needs_vlm"] = True
            elif shape.shape_type == MSO.GROUP:
                # Group shapes are not decomposed (too fiddly — child structure and alignment break easily): leave a note in place and hand the whole slide to the VLM.
                try:
                    n = len(shape.shapes)
                except Exception:
                    n = "?"
                rec["lines"] = [_tag([*_shape_bits(shape, with_size=True), f"group-of-{n}-shapes", "not-decomposed", "vlm"])]
                rec["needs_vlm"] = True
            elif _is_line(shape):
                rec["kind"] = "arrow"
                rec["info"] = _arrow_info(shape, idx, shapes, bbox)
            elif shape.has_text_frame and shape.text_frame.text.strip():
                fmt_bits, lines = _pptx_text_md(shape)
                box = text_to_box.get(idx)
                meta = _shape_bits(shapes[box], with_size=True) if box is not None else _shape_bits(shape)
                pos = meta[0] if meta and meta[0].startswith("@") else ""
                rest = [b for b in meta if not b.startswith("@")]
                rec.update(kind="text", meta=meta, fmt=fmt_bits, lines=lines,
                           pos=pos, sig=tuple(rest + fmt_bits), label=_shape_text(shape))
            else:  # purely visual element with no text
                rec["lines"] = [_tag([*_shape_bits(shape, with_size=True), "vlm"])]
                rec["needs_vlm"] = True
        except Exception as e:
            rec["lines"] = [f"`shape parse error: {type(e).__name__}: {e}`"]
        records.append(rec)

    # Clustering: members are connected into a group (union-find) when they pairwise share
    # ">= _GROUP_MIN_COMMON of the backtick bits (size / geometry / fill / font size / colour)";
    # which specific parameters is not prescribed — any 3 or more identical ones count.
    for r in records:
        if r["kind"] == "text":
            r["nonpos"] = [b for b in r["meta"] if not b.startswith("@")] + r["fmt"]
            r["bitset"] = set(r["nonpos"])
    text_recs = [r for r in records if r["kind"] == "text"]
    parent = {r["idx"]: r["idx"] for r in text_recs}

    def _find(x):
        while parent[x] != x:
            parent[x] = parent[parent[x]]
            x = parent[x]
        return x

    for a in range(len(text_recs)):
        for b in range(a + 1, len(text_recs)):
            if len(text_recs[a]["bitset"] & text_recs[b]["bitset"]) >= _GROUP_MIN_COMMON:
                ra, rb = _find(text_recs[a]["idx"]), _find(text_recs[b]["idx"])
                if ra != rb:
                    parent[ra] = rb
    comps = defaultdict(list)
    for r in text_recs:
        comps[_find(r["idx"])].append(r)
    # A valid group: >= 2 members whose overall intersection still has >= _GROUP_MIN_COMMON parameters (stops chained merges producing a weak intersection)
    rec_group, group_members = {}, {}
    for root, members in comps.items():
        if len(members) < 2:
            continue
        if len(set.intersection(*[m["bitset"] for m in members])) < _GROUP_MIN_COMMON:
            continue
        group_members[root] = members
        for m in members:
            rec_group[m["idx"]] = root

    label_to_group = {}
    for root, members in group_members.items():
        for m in members:
            if m["label"]:
                label_to_group[m["label"][:24]] = root

    # Arrows: both ends in one group → fold into that group's flow; otherwise standalone
    group_edges = defaultdict(list)
    for r in records:
        if r["kind"] == "arrow":
            s, d = r["info"]["src"], r["info"]["dst"]
            if s and d and label_to_group.get(s) and label_to_group.get(s) == label_to_group.get(d):
                group_edges[label_to_group[s]].append((s, d))
                r["consumed"] = True
            elif not r["info"]["connected"]:
                r["needs_vlm"] = True

    # Content is emitted in reading order; purely visual elements needing a VLM (pictures / groups / visuals / unconnected arrows) are gathered at the end of the slide rather than interleaved
    out, emitted, vlm = [], set(), []
    for r in records:
        if r.get("consumed"):
            continue
        if r.get("needs_vlm"):
            vlm.extend([r["info"]["line"]] if r["kind"] == "arrow" else r["lines"])
            continue
        out.append("")
        if r["kind"] == "text" and r["idx"] in rec_group:
            root = rec_group[r["idx"]]
            if root in emitted:
                out.pop()
                continue
            emitted.add(root)
            members = group_members[root]
            # Shared parameters = the intersection of all members' bitsets (ordered by the first member); hoisted to the group header
            common = [b for b in members[0]["nonpos"]
                      if all(b in mm["bitset"] for mm in members)]
            out.append(_tag([f"group ×{len(members)}", *common]))
            for m in members:
                extra = [m["pos"]] + [b for b in m["nonpos"] if b not in common]
                mtag = _tag([b for b in extra if b])  # position + this member's own parameters (e.g. its individual fill)
                if len(m["lines"]) == 1:
                    out.append(f"- {mtag} {m['lines'][0]}".strip())
                else:
                    out.append(f"- {mtag}".rstrip())
                    out.extend("  " + ln for ln in m["lines"])
            flow = _order_flow(group_edges.get(root, []))
            if flow:
                out.append(f"- flow: {flow}")
        elif r["kind"] == "text":
            tag = _tag(r["meta"] + r["fmt"])
            if len(r["lines"]) == 1 and tag:
                out.append(f"{tag} {r['lines'][0]}")
            else:
                if tag:
                    out.append(tag)
                out.extend(r["lines"])
        elif r["kind"] == "arrow":
            out.append(r["info"]["line"])
        else:
            out.extend(r["lines"])
    if vlm:
        out += ["", "▸ visual elements (need VLM)", *vlm]
    return out, bool(vlm)


def _pptx_to_md(path: str) -> str:
    """pptx → markdown, slide by slide and shape by shape in **reading order (XML order)**.
    One compact `…` meta per shape (only signal that
    deviates from the defaults: position / placeholder / non-rectangular geometry / fill /
    font size and colour / chart / vlm), then the content. `…` = parser-added, not original.
    Density: defaults like black and rectangle are not annotated; a shape's shared font
    size / colour is hoisted into meta once and runs record only deviations.
    Structure recovery: "background colour block + the text label inside it" merge into one
    line; lines / arrows get a direction and the labels they connect."""
    _ensure("pptx", "python-pptx")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    out = [
        f"<!-- pptx readout · slide {_emu_in(prs.slide_width)}×{_emu_in(prs.slide_height)}in."
        " `…` (inline code) = parser-added metadata, not slide content: @x,y = position"
        " in inches; NNpt/#hex = font size/color (deviations only); <u>underline</u>;"
        " arrow→/connects = connector; vlm = purely visual element (image/grouping/"
        "overlap), needs a VLM. -->",
        "",
    ]
    vlm_pages = []  # slide numbers containing purely visual elements (deduped; the whole slide goes to the VLM instead of listing shape by shape)
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"<!-- slide {i} -->")
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                out.append("")
                out.append("> **Notes:** " + note.replace("\n", "\n> "))
        shapes = list(slide.shapes)
        bbox = [_bbox(s) for s in shapes]
        text_to_box, used = _pair_label_boxes(shapes, bbox)
        lines, needs_vlm = _pptx_slide_lines(shapes, bbox, text_to_box, used, MSO_SHAPE_TYPE)
        out.extend(lines)
        out.append("")
        if needs_vlm:
            vlm_pages.append(i)
    if vlm_pages:
        out.append("<!-- needs VLM -->")
        out.append("Slides: " + ", ".join(str(p) for p in vlm_pages))
        out.append("The above slides contain purely visual elements that the text layer can't capture. Use VLM to review these slides if necessary.")
    return "\n".join(out)
